"""PowerPoint (.pptx) loader using python-pptx. Each slide becomes a page."""
from __future__ import annotations
import uuid
from pathlib import Path

from raglib.loaders.base import BaseLoader
from raglib.models.document import Document, ImageRef


class PptxLoader(BaseLoader):
    """
    Loads .pptx files. Each slide becomes a page_map entry.
    Extracts: title, body text, speaker notes, and embedded images.
    """

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    def load(self, path: str) -> list[Document]:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches  # type: ignore  # noqa: F401
        except ImportError:
            raise ImportError(
                "python-pptx is required to load .pptx files. "
                "Install it with: pip install python-pptx"
            )

        p = Path(path)
        prs = Presentation(str(p))
        doc_id = p.stem + "_" + uuid.uuid4().hex[:8]

        text_parts: list[str] = []
        page_map: list[dict] = []
        image_refs: list[ImageRef] = []
        char_offset = 0
        image_index = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_title = ""
            slide_body_lines: list[str] = []
            slide_notes = ""

            # Extract title and body text from shapes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    # Check for images
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                img_bytes = shape.image.blob
                                mime = shape.image.content_type or "image/png"
                                image_refs.append(ImageRef(
                                    image_id=f"{p.stem}_s{slide_num}_img{image_index}",
                                    char_offset=char_offset,
                                    image_bytes=img_bytes,
                                    mime_type=mime,
                                    page_num=slide_num,
                                ))
                                image_index += 1
                            except Exception:
                                pass
                    except Exception:
                        pass
                    continue

                text_frame = shape.text_frame
                shape_text = text_frame.text.strip()

                if not shape_text:
                    continue

                # Detect title placeholder
                try:
                    from pptx.util import Pt  # type: ignore  # noqa: F401
                    ph = shape.placeholder_format
                    if ph is not None and ph.idx == 0:
                        slide_title = shape_text
                        continue
                except Exception:
                    pass

                slide_body_lines.append(shape_text)

            # Extract speaker notes
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_notes = notes_text
            except Exception:
                pass

            # Build slide text block
            slide_parts = []
            if slide_title:
                slide_parts.append(f"# {slide_title}")
            slide_parts.extend(slide_body_lines)
            if slide_notes:
                slide_parts.append(f"\n[Notes]: {slide_notes}")

            slide_text = "\n".join(slide_parts)
            separator = "\n\n" if text_parts else ""
            full_part = separator + slide_text

            start_char = char_offset + len(separator)
            end_char = char_offset + len(full_part)

            page_map.append({
                "page_num": slide_num,
                "start_char": start_char,
                "end_char": end_char,
                "heading": slide_title or f"Slide {slide_num}",
            })

            text_parts.append(full_part)
            char_offset = end_char

        full_text = "".join(text_parts)

        return [
            Document(
                doc_id=doc_id,
                text=full_text,
                source_path=str(p.resolve()),
                source_name=p.name,
                loader="PptxLoader",
                file_extension=p.suffix.lower(),
                file_size_bytes=p.stat().st_size,
                page_count=len(page_map),
                page_map=page_map,
                image_refs=image_refs,
            )
        ]
